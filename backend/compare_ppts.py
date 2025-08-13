import os
import dotenv
from pptx import Presentation
from openai import OpenAI
dotenv.load_dotenv()


# --- Configuration ---
# Make sure to set your OpenAI API key as an environment variable
client = OpenAI()
PPT_DIRECTORY = "PPT"
PPT_FILE_1 = "慢性淋巴细胞白血病最新研究进展_1-3月.pptx"
PPT_FILE_2 = "慢性淋巴细胞白血病最新研究进展_4-6月.pptx"
# --- End Configuration ---

def extract_text_from_ppt(ppt_path):
    """Extracts all text from a PowerPoint file."""
    try:
        prs = Presentation(ppt_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)
    except Exception as e:
        return f"Error reading {os.path.basename(ppt_path)}: {e}"

def compare_ppt_content(text1, text2, filename1, filename2):
    """Compares the text content of two PPTs using OpenAI's API."""
    prompt = f"""
    以下是两个 PowerPoint 演示文稿的文本内容。  
    请仔细阅读并对比它们的内容，提炼并总结它们的主要差异，要求用中文输出。  
    
    文件: {filename1}  
    ---  
    {text1}  
    ---  
    
    文件: {filename2}  
    ---  
    {text2}  
    ---  
    
    请用简洁且结构化的方式，按主题或要点列出它们的关键差异。
    """

    response = client.chat.completions.create(
        model="gpt-4o",  # Or a more advanced model like gpt-4
        messages=[
            {"role": "system", "content": "You are a helpful assistant that summarizes differences in documents."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content

def main():
    """Main function to extract text and compare two PPTs."""
    ppt1_path = os.path.join(PPT_DIRECTORY, PPT_FILE_1)
    ppt2_path = os.path.join(PPT_DIRECTORY, PPT_FILE_2)

    # Check if files exist
    if not os.path.exists(ppt1_path):
        print(f"Error: File not found at {ppt1_path}")
        return
    if not os.path.exists(ppt2_path):
        print(f"Error: File not found at {ppt2_path}")
        return

    print("Extracting text from PowerPoint files...")
    ppt1_text = extract_text_from_ppt(ppt1_path)
    ppt2_text = extract_text_from_ppt(ppt2_path)

    if "Error reading" in ppt1_text:
        print(ppt1_text)
        return
    if "Error reading" in ppt2_text:
        print(ppt2_text)
        return

    print("Comparing PowerPoint content using OpenAI...")
    comparison_summary = compare_ppt_content(ppt1_text, ppt2_text, PPT_FILE_1, PPT_FILE_2)

    print("\n--- Comparison Summary ---")
    print(comparison_summary)
    print("--------------------------\n")

if __name__ == "__main__":
    main()
